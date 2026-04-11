"""Builder stage - PPTX assembly with python-pptx."""

from __future__ import annotations

import io
import logging
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from noteeditor.models.content import ExtractedImage, FontMatch, OCRResult, TextStyle
from noteeditor.models.layout import LayoutRegion, LayoutResult, RegionLabel
from noteeditor.models.page import PageImage
from noteeditor.models.slide import ImageBlock, SlideContent, TextBlock

logger = logging.getLogger(__name__)

# Standard aspect ratio → slide dimensions in inches
# (width, height)
_STANDARD_RATIOS: list[tuple[float, float, float]] = [
    (16 / 9, 13.333, 7.5),
    (4 / 3, 10.0, 7.5),
    (16 / 10, 10.0, 6.25),
]

_EMU_PER_INCH = 914400


def detect_slide_dimensions(aspect_ratio: float) -> tuple[int, int]:
    """Map a page aspect ratio to PPTX slide dimensions in EMU.

    Finds the closest standard ratio (by absolute difference) and
    returns the corresponding slide dimensions.

    Args:
        aspect_ratio: Page width / height (e.g., 1.778 for 16:9).

    Returns:
        Tuple of (width_emu, height_emu).
    """
    best_diff = float("inf")
    best_width_in = 10.0
    best_height_in = 7.5

    for ratio, width_in, height_in in _STANDARD_RATIOS:
        diff = abs(aspect_ratio - ratio)
        if diff < best_diff:
            best_diff = diff
            best_width_in = width_in
            best_height_in = height_in

    return int(best_width_in * _EMU_PER_INCH), int(best_height_in * _EMU_PER_INCH)


def _image_to_bytes(image: np.ndarray) -> bytes:
    """Convert an RGB numpy array to PNG bytes."""
    pil_image = Image.fromarray(image)
    buf = io.BytesIO()
    pil_image.save(buf, format="PNG")
    return buf.getvalue()


def build_pptx(
    pages: tuple[PageImage, ...],
    output_path: Path,
) -> Path:
    """Build a PPTX file from page images (screenshot mode).

    Each page's rendered image is added as a full-slide background picture.
    Slide dimensions are determined by the first page's aspect ratio.

    Args:
        pages: Tuple of PageImage objects to include as slides.
        output_path: Where to write the .pptx file.

    Returns:
        The output_path (for chaining).
    """
    # Determine slide dimensions from first page
    if pages:
        width_emu, height_emu = detect_slide_dimensions(pages[0].aspect_ratio)
    else:
        # Default to 16:9 for empty presentations
        width_emu, height_emu = detect_slide_dimensions(16 / 9)

    prs = Presentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        image_bytes = _image_to_bytes(page.image)
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            Emu(0),
            Emu(0),
            Emu(width_emu),
            Emu(height_emu),
        )

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    logger.info("Built PPTX: %s (%d slides)", output_path, len(pages))
    return output_path


# --- Editable mode (Feature 011) ---


def _make_fallback_font_match(region_id: str, label: RegionLabel) -> FontMatch:
    """Create a fallback FontMatch for v0.2.0 MVP (Arial).

    Real font matching is deferred to v0.4.0.
    """
    return FontMatch(
        region_id=region_id,
        label=label,
        font_name="Arial",
        font_path=None,
        system_fallback="Arial",
        is_fallback=True,
    )


def _estimate_font_size(bbox_height_px: float, dpi: int) -> int:
    """Estimate font size in points from bbox height in pixels.

    Formula: font_size_pt = bbox_height_px * (72 / dpi) * 0.8
    The 0.8 factor accounts for typical line-height vs font-size ratio.
    """
    return max(1, int(bbox_height_px * (72 / dpi) * 0.8))


def assemble_slide(
    page_image: PageImage,
    layout_result: LayoutResult,
    ocr_results: tuple[OCRResult, ...],
    background_image: np.ndarray | None = None,
    image_results: tuple[ExtractedImage, ...] = (),
    font_matches: tuple[FontMatch, ...] = (),
    style_results: tuple[TextStyle, ...] = (),
) -> SlideContent:
    """Assemble OCR results into SlideContent for the builder.

    Maps OCRResults to TextBlocks and ExtractedImages to ImageBlocks
    using LayoutRegion data for positioning.

    Args:
        page_image: Source page image.
        layout_result: Layout detection result with region positions.
        ocr_results: OCR extraction results with text content.
        background_image: Clean background with text removed (v0.3.0).
        image_results: Extracted images for IMAGE regions (Feature 014).
        font_matches: Font matching results for text regions (Feature 015).
        style_results: Estimated text styles for text regions (Feature 016).

    Returns:
        Frozen SlideContent ready for build_editable_pptx().
    """
    region_by_id: dict[str, LayoutRegion] = {
        r.region_id: r for r in layout_result.regions
    }
    font_by_region: dict[str, FontMatch] = {
        f.region_id: f for f in font_matches
    }
    style_by_region: dict[str, TextStyle] = {
        s.region_id: s for s in style_results
    }

    text_blocks: list[TextBlock] = []
    for ocr in ocr_results:
        region = region_by_id.get(ocr.region_id)
        if region is None:
            continue

        font_match = font_by_region.get(ocr.region_id)
        if font_match is None:
            font_match = _make_fallback_font_match(ocr.region_id, region.label)
        style = style_by_region.get(ocr.region_id)
        text_blocks.append(
            TextBlock(
                region_id=ocr.region_id,
                bbox=region.bbox,
                text=ocr.text,
                font_match=font_match,
                is_formula=ocr.is_formula,
                formula_latex=ocr.formula_latex,
                style=style,
            ),
        )

    image_blocks: list[ImageBlock] = []
    for img in image_results:
        image_blocks.append(
            ImageBlock(
                region_id=img.region_id,
                bbox=img.bbox,
                image=img.image,
                source=img.source,
            ),
        )

    return SlideContent(
        page_number=page_image.page_number,
        background_image=background_image,
        full_page_image=page_image.image,
        text_blocks=tuple(text_blocks),
        image_blocks=tuple(image_blocks),
        status="success",
    )


def _add_text_box(
    slide: Presentation.slides,
    text_block: TextBlock,
    dpi: int,
    has_clean_background: bool = False,
) -> None:
    """Add a text box to a slide.

    When has_clean_background is True (background extraction produced a clean
    image), the text box has no fill (transparent). Otherwise, white fill is
    used to cover the original text in the screenshot background.

    Position and size are converted from pixel coordinates to EMU using DPI.
    Alignment is set based on the region label (TITLE=center, others=left).
    """
    bbox = text_block.bbox
    emu_per_px = _EMU_PER_INCH / dpi

    left = int(bbox.x * emu_per_px)
    top = int(bbox.y * emu_per_px)
    width = int(bbox.width * emu_per_px)
    height = int(bbox.height * emu_per_px)

    textbox = slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height),
    )

    if has_clean_background:
        # Transparent fill — text sits on the already-cleaned background
        textbox.fill.background()
    else:
        # White fill to cover original text in the screenshot
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Set text content and formatting
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text_block.text

    run = p.runs[0]

    # Use style if available, otherwise fall back to estimation
    if text_block.style is not None:
        run.font.size = Pt(text_block.style.font_size_pt)
        r, g, b = text_block.style.font_color_rgb
        run.font.color.rgb = RGBColor(r, g, b)
    else:
        run.font.size = Pt(_estimate_font_size(bbox.height, dpi))

    run.font.name = text_block.font_match.font_name

    # Alignment based on region type
    if text_block.font_match.label == RegionLabel.TITLE:
        p.alignment = PP_ALIGN.CENTER
    else:
        p.alignment = PP_ALIGN.LEFT


def _add_image_block(
    slide: Presentation.slides,
    image_block: ImageBlock,
    dpi: int,
) -> None:
    """Add an image block to a slide.

    Position and size are converted from pixel coordinates to EMU using DPI.
    """
    bbox = image_block.bbox
    emu_per_px = _EMU_PER_INCH / dpi

    left = int(bbox.x * emu_per_px)
    top = int(bbox.y * emu_per_px)
    width = int(bbox.width * emu_per_px)
    height = int(bbox.height * emu_per_px)

    image_bytes = _image_to_bytes(image_block.image)
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        Emu(left), Emu(top), Emu(width), Emu(height),
    )


def _set_slide_background_image(slide: object, image_bytes: bytes) -> None:
    """Set a native slide background image via XML.

    Uses blipFill on the slide's bgPr element so the background is
    not selectable or movable by the user (unlike add_picture).
    """
    slide_part = slide.part  # type: ignore[union-attr]
    image_stream = io.BytesIO(image_bytes)
    _image_part, rId = slide_part.get_or_add_image_part(image_stream)

    bg_elem = slide.background._element  # type: ignore[union-attr]

    # Remove existing bgPr if present
    existing = bg_elem.find(qn("p:bgPr"))
    if existing is not None:
        bg_elem.remove(existing)

    bgPr = parse_xml(
        '<p:bgPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )
    blipFill = parse_xml(
        '<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<a:blip r:embed="{rId}"/>'
        "<a:stretch><a:fillRect/></a:stretch>"
        "</a:blipFill>"
    )
    bgPr.append(blipFill)
    bg_elem.append(bgPr)


def build_editable_pptx(
    pages: tuple[SlideContent, ...],
    output_path: Path,
    dpi: int = 300,
) -> Path:
    """Build an editable PPTX with text boxes over backgrounds.

    When a clean background image is available, it is set as the native
    slide background (not selectable/movable) and text boxes are transparent.
    When no clean background exists, the full page screenshot is placed as
    an add_picture fallback with white-filled text boxes.

    Args:
        pages: Tuple of SlideContent objects to render as slides.
        output_path: Where to write the .pptx file.
        dpi: DPI used for pixel-to-EMU coordinate conversion.

    Returns:
        The output_path (for chaining).
    """
    if not pages:
        prs = Presentation()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    # Determine slide dimensions from first page's image
    first_h, first_w = pages[0].full_page_image.shape[:2]
    aspect_ratio = first_w / first_h
    width_emu, height_emu = detect_slide_dimensions(aspect_ratio)

    prs = Presentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)

    blank_layout = prs.slide_layouts[6]

    for slide_content in pages:
        slide = prs.slides.add_slide(blank_layout)

        has_clean_bg = slide_content.background_image is not None

        if has_clean_bg:
            # Native slide background (not selectable/movable)
            image_bytes = _image_to_bytes(slide_content.background_image)
            _set_slide_background_image(slide, image_bytes)
        else:
            # Fallback: add_picture as full-slide image
            image_bytes = _image_to_bytes(slide_content.full_page_image)
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Emu(0), Emu(0),
                Emu(width_emu), Emu(height_emu),
            )

        # Text boxes and images only for successful pages (not fallback)
        if slide_content.status != "fallback":
            for text_block in slide_content.text_blocks:
                _add_text_box(slide, text_block, dpi, has_clean_bg)
            for image_block in slide_content.image_blocks:
                _add_image_block(slide, image_block, dpi)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    logger.info(
        "Built editable PPTX: %s (%d slides)", output_path, len(pages),
    )
    return output_path
