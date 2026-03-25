"""Tests for stages/builder.py - PPTX building."""

from __future__ import annotations

from pathlib import Path

import numpy as np

from noteeditor.models.page import PageImage
from noteeditor.stages.builder import build_pptx, detect_slide_dimensions


def _make_page_image(
    page_number: int = 0,
    width_px: int = 4000,
    height_px: int = 2250,
    dpi: int = 300,
) -> PageImage:
    """Create a PageImage with synthetic image data."""
    image = np.zeros((height_px, width_px, 3), dtype=np.uint8)
    aspect_ratio = width_px / height_px if height_px > 0 else 1.0
    return PageImage(
        page_number=page_number,
        width_px=width_px,
        height_px=height_px,
        dpi=dpi,
        aspect_ratio=aspect_ratio,
        image=image,
    )


class TestDetectSlideDimensions:
    """Tests for detect_slide_dimensions."""

    def test_16_9_ratio(self) -> None:
        """16:9 aspect ratio maps to standard widescreen dimensions."""
        width_emu, height_emu = detect_slide_dimensions(16 / 9)
        # 13.333in × 7.5in in EMU (914400 EMU/inch)
        expected_width = int(13.333 * 914400)
        expected_height = int(7.5 * 914400)
        assert width_emu == expected_width
        assert height_emu == expected_height

    def test_4_3_ratio(self) -> None:
        """4:3 aspect ratio maps to standard dimensions."""
        width_emu, height_emu = detect_slide_dimensions(4 / 3)
        # 10in × 7.5in in EMU
        expected_width = int(10 * 914400)
        expected_height = int(7.5 * 914400)
        assert width_emu == expected_width
        assert height_emu == expected_height

    def test_16_10_ratio(self) -> None:
        """16:10 aspect ratio maps to closest standard ratio."""
        width_emu, height_emu = detect_slide_dimensions(16 / 10)
        # 16:10 ≈ 1.6, closest to 16:9 (1.778) or 4:3 (1.333)
        assert width_emu > 0
        assert height_emu > 0

    def test_1_1_ratio(self) -> None:
        """1:1 aspect ratio returns valid dimensions."""
        width_emu, height_emu = detect_slide_dimensions(1.0)
        assert width_emu > 0
        assert height_emu > 0

    def test_exact_16_9_float(self) -> None:
        """Exact 16:9 float (1.777...) matches widescreen."""
        ratio = 16 / 9  # 1.7777...
        width_emu, height_emu = detect_slide_dimensions(ratio)
        expected_width = int(13.333 * 914400)
        expected_height = int(7.5 * 914400)
        assert width_emu == expected_width
        assert height_emu == expected_height

    def test_returns_tuple_of_ints(self) -> None:
        """Returns a tuple of two positive integers."""
        result = detect_slide_dimensions(16 / 9)
        assert isinstance(result, tuple)
        assert len(result) == 2
        assert all(isinstance(v, int) for v in result)
        assert all(v > 0 for v in result)


class TestBuildPptx:
    """Tests for build_pptx."""

    def test_single_page_creates_pptx(self, tmp_path: Path) -> None:
        """Single page produces a valid PPTX file."""
        pages = (_make_page_image(page_number=0),)
        output = tmp_path / "output.pptx"

        result = build_pptx(pages, output)

        assert result == output
        assert output.exists()
        assert output.stat().st_size > 0

    def test_multi_page_creates_pptx(self, tmp_path: Path) -> None:
        """Multiple pages produce a PPTX with correct slide count."""
        pages = (
            _make_page_image(page_number=0),
            _make_page_image(page_number=1),
            _make_page_image(page_number=2),
        )
        output = tmp_path / "output.pptx"

        build_pptx(pages, output)

        from pptx import Presentation

        prs = Presentation(str(output))
        assert len(prs.slides) == 3

    def test_empty_pages_creates_empty_pptx(self, tmp_path: Path) -> None:
        """Empty page tuple creates a PPTX with zero slides."""
        output = tmp_path / "empty.pptx"

        result = build_pptx((), output)

        assert result == output
        assert output.exists()
        from pptx import Presentation

        prs = Presentation(str(output))
        assert len(prs.slides) == 0

    def test_slide_dimensions_match_aspect_ratio(self, tmp_path: Path) -> None:
        """PPTX slide dimensions match the first page's aspect ratio."""
        # 16:9 ratio: 4000/2250 ≈ 1.778
        pages = (_make_page_image(width_px=4000, height_px=2250),)
        output = tmp_path / "output.pptx"

        build_pptx(pages, output)

        from pptx import Presentation

        prs = Presentation(str(output))
        expected_width, expected_height = detect_slide_dimensions(4000 / 2250)
        assert prs.slide_width == expected_width
        assert prs.slide_height == expected_height

    def test_output_path_parent_created(self, tmp_path: Path) -> None:
        """Parent directories are created if they don't exist."""
        output = tmp_path / "nested" / "dir" / "output.pptx"
        pages = (_make_page_image(),)

        build_pptx(pages, output)

        assert output.exists()

    def test_4_3_aspect_ratio_dimensions(self, tmp_path: Path) -> None:
        """4:3 pages produce correct PPTX dimensions."""
        # 1024x768 ≈ 4:3
        pages = (_make_page_image(width_px=1024, height_px=768),)
        output = tmp_path / "output.pptx"

        build_pptx(pages, output)

        from pptx import Presentation

        prs = Presentation(str(output))
        expected_width, expected_height = detect_slide_dimensions(1024 / 768)
        assert prs.slide_width == expected_width
        assert prs.slide_height == expected_height

    def test_slides_contain_images(self, tmp_path: Path) -> None:
        """Each slide has at least one image shape."""
        pages = (_make_page_image(page_number=0), _make_page_image(page_number=1))
        output = tmp_path / "output.pptx"

        build_pptx(pages, output)

        from pptx import Presentation

        prs = Presentation(str(output))
        for slide in prs.slides:
            image_shapes = [
                s for s in slide.shapes if s.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            ]
            assert len(image_shapes) >= 1
