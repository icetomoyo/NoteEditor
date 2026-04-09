"""Tests for stages/builder.py - Editable PPTX building (Feature 011)."""

from __future__ import annotations

from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from noteeditor.models.content import ExtractedImage, FontMatch, OCRResult, TextStyle
from noteeditor.models.layout import BoundingBox, LayoutRegion, LayoutResult, RegionLabel
from noteeditor.models.page import PageImage
from noteeditor.models.slide import ImageBlock, SlideContent, TextBlock
from noteeditor.stages.builder import (
    _add_text_box,
    _estimate_font_size,
    _make_fallback_font_match,
    assemble_slide,
    build_editable_pptx,
)


def _make_page_image(
    page_number: int = 0,
    width_px: int = 4000,
    height_px: int = 2250,
    dpi: int = 300,
) -> PageImage:
    """Create a PageImage with synthetic image data."""
    image = np.random.randint(0, 255, (height_px, width_px, 3), dtype=np.uint8)
    aspect_ratio = width_px / height_px if height_px > 0 else 1.0
    return PageImage(
        page_number=page_number,
        width_px=width_px,
        height_px=height_px,
        dpi=dpi,
        aspect_ratio=aspect_ratio,
        image=image,
    )


def _make_layout_result(
    page_number: int = 0,
    regions: tuple[LayoutRegion, ...] = (),
) -> LayoutResult:
    """Create a LayoutResult with given regions."""
    return LayoutResult(page_number=page_number, regions=regions)


def _make_layout_region(
    region_id: str = "page0_region0",
    label: RegionLabel = RegionLabel.BODY_TEXT,
    confidence: float = 0.9,
    x: float = 100.0,
    y: float = 200.0,
    width: float = 500.0,
    height: float = 60.0,
) -> LayoutRegion:
    """Create a LayoutRegion with given parameters."""
    return LayoutRegion(
        bbox=BoundingBox(x=x, y=y, width=width, height=height),
        label=label,
        confidence=confidence,
        region_id=region_id,
    )


def _make_ocr_result(
    region_id: str = "page0_region0",
    text: str = "Hello World",
    confidence: float = 0.95,
    is_formula: bool = False,
    formula_latex: str | None = None,
) -> OCRResult:
    """Create an OCRResult with given parameters."""
    return OCRResult(
        region_id=region_id,
        text=text,
        confidence=confidence,
        is_formula=is_formula,
        formula_latex=formula_latex,
    )


def _make_text_block(
    region_id: str = "page0_region0",
    text: str = "Hello",
    label: RegionLabel = RegionLabel.BODY_TEXT,
    x: float = 100.0,
    y: float = 200.0,
    width: float = 500.0,
    height: float = 60.0,
) -> TextBlock:
    """Create a TextBlock with fallback FontMatch."""
    font_match = _make_fallback_font_match(region_id, label)
    return TextBlock(
        region_id=region_id,
        bbox=BoundingBox(x=x, y=y, width=width, height=height),
        text=text,
        font_match=font_match,
        is_formula=False,
    )


def _make_slide_content(
    page_number: int = 0,
    text_blocks: tuple[TextBlock, ...] = (),
    image_blocks: tuple[ImageBlock, ...] = (),
    status: str = "success",
    width_px: int = 4000,
    height_px: int = 2250,
    background_image: np.ndarray | None = None,
) -> SlideContent:
    """Create a SlideContent with synthetic data."""
    image = np.random.randint(0, 255, (height_px, width_px, 3), dtype=np.uint8)
    return SlideContent(
        page_number=page_number,
        background_image=background_image,
        full_page_image=image,
        text_blocks=text_blocks,
        image_blocks=image_blocks,
        status=status,  # type: ignore[arg-type]
    )


class TestEstimateFontSize:
    """Tests for _estimate_font_size()."""

    def test_basic_estimation(self) -> None:
        """Standard bbox height at 300 DPI produces correct font size."""
        # height=60px, dpi=300 → 60 * 72/300 * 0.8 = 11.52 → 11
        result = _estimate_font_size(60.0, 300)
        assert result == 11

    def test_returns_integer(self) -> None:
        """Result is always an integer."""
        result = _estimate_font_size(100.0, 300)
        assert isinstance(result, int)

    def test_minimum_font_size(self) -> None:
        """Very small bbox returns minimum font size of 1."""
        result = _estimate_font_size(0.5, 300)
        assert result >= 1

    def test_zero_height_returns_minimum(self) -> None:
        """Zero height bbox returns minimum font size of 1."""
        result = _estimate_font_size(0.0, 300)
        assert result >= 1

    def test_higher_dpi_smaller_font(self) -> None:
        """Higher DPI produces smaller font for the same pixel height."""
        at_150 = _estimate_font_size(60.0, 150)
        at_300 = _estimate_font_size(60.0, 300)
        assert at_150 > at_300

    def test_large_title(self) -> None:
        """Large title bbox produces proportionally large font."""
        # height=120px, dpi=300 → 120 * 72/300 * 0.8 = 23.04 → 23
        result = _estimate_font_size(120.0, 300)
        assert result == 23


class TestMakeFallbackFontMatch:
    """Tests for _make_fallback_font_match()."""

    def test_creates_font_match(self) -> None:
        """Returns a FontMatch instance."""
        result = _make_fallback_font_match("r1", RegionLabel.BODY_TEXT)
        assert isinstance(result, FontMatch)

    def test_is_fallback(self) -> None:
        """Font match is marked as fallback."""
        result = _make_fallback_font_match("r1", RegionLabel.BODY_TEXT)
        assert result.is_fallback is True

    def test_uses_arial(self) -> None:
        """Fallback font is Arial."""
        result = _make_fallback_font_match("r1", RegionLabel.BODY_TEXT)
        assert result.font_name == "Arial"
        assert result.system_fallback == "Arial"

    def test_preserves_region_id(self) -> None:
        """Region ID is preserved."""
        result = _make_fallback_font_match("page0_region3", RegionLabel.TITLE)
        assert result.region_id == "page0_region3"

    def test_preserves_label(self) -> None:
        """Region label is preserved."""
        result = _make_fallback_font_match("r1", RegionLabel.TITLE)
        assert result.label == RegionLabel.TITLE

    def test_no_font_path(self) -> None:
        """No specific font path (system font)."""
        result = _make_fallback_font_match("r1", RegionLabel.BODY_TEXT)
        assert result.font_path is None


class TestAssembleSlide:
    """Tests for assemble_slide()."""

    def test_basic_assembly(self) -> None:
        """Creates SlideContent from PageImage + LayoutResult + OCRResults."""
        page = _make_page_image()
        region = _make_layout_region()
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result()

        result = assemble_slide(page, layout, (ocr,))

        assert isinstance(result, SlideContent)
        assert result.page_number == 0
        assert len(result.text_blocks) == 1

    def test_maps_ocr_to_text_blocks(self) -> None:
        """Each OCRResult becomes a TextBlock."""
        r1 = _make_layout_region(region_id="r1")
        r2 = _make_layout_region(region_id="r2", x=100, y=300)
        layout = _make_layout_result(regions=(r1, r2))

        o1 = _make_ocr_result(region_id="r1", text="Hello")
        o2 = _make_ocr_result(region_id="r2", text="World")

        page = _make_page_image()
        result = assemble_slide(page, layout, (o1, o2))

        assert len(result.text_blocks) == 2
        texts = {tb.text for tb in result.text_blocks}
        assert texts == {"Hello", "World"}

    def test_text_block_has_fallback_font(self) -> None:
        """TextBlock uses fallback FontMatch."""
        page = _make_page_image()
        region = _make_layout_region()
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result()

        result = assemble_slide(page, layout, (ocr,))

        assert result.text_blocks[0].font_match.is_fallback is True
        assert result.text_blocks[0].font_match.font_name == "Arial"

    def test_preserves_bbox_from_layout_region(self) -> None:
        """TextBlock bbox comes from LayoutRegion, not OCR."""
        page = _make_page_image()
        region = _make_layout_region(x=50, y=100, width=600, height=80)
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result()

        result = assemble_slide(page, layout, (ocr,))

        bbox = result.text_blocks[0].bbox
        assert bbox.x == 50.0
        assert bbox.y == 100.0
        assert bbox.width == 600.0
        assert bbox.height == 80.0

    def test_preserves_page_number(self) -> None:
        """Page number from PageImage is preserved."""
        page = _make_page_image(page_number=5)
        layout = _make_layout_result(page_number=5)

        result = assemble_slide(page, layout, ())

        assert result.page_number == 5

    def test_empty_ocr_results(self) -> None:
        """Empty OCR results produce empty text_blocks."""
        page = _make_page_image()
        layout = _make_layout_result(regions=(_make_layout_region(),))

        result = assemble_slide(page, layout, ())

        assert result.text_blocks == ()

    def test_status_is_success(self) -> None:
        """Successful assembly has status='success'."""
        page = _make_page_image()
        layout = _make_layout_result()

        result = assemble_slide(page, layout, ())

        assert result.status == "success"

    def test_background_image_is_none_by_default(self) -> None:
        """Without background_image param, field is None."""
        page = _make_page_image()
        layout = _make_layout_result()

        result = assemble_slide(page, layout, ())

        assert result.background_image is None

    def test_background_image_set_when_provided(self) -> None:
        """When background_image is provided, it is set in SlideContent."""
        page = _make_page_image()
        layout = _make_layout_result()
        bg = np.ones((2250, 4000, 3), dtype=np.uint8) * 200

        result = assemble_slide(page, layout, (), background_image=bg)

        assert result.background_image is bg

    def test_image_blocks_empty_when_no_results(self) -> None:
        """Without image_results, image_blocks is empty."""
        page = _make_page_image()
        layout = _make_layout_result()

        result = assemble_slide(page, layout, ())

        assert result.image_blocks == ()

    def test_image_blocks_populated_from_results(self) -> None:
        """image_results are converted to ImageBlocks."""
        page = _make_page_image()
        layout = _make_layout_result()
        img = np.zeros((100, 200, 3), dtype=np.uint8)
        extracted = ExtractedImage(
            region_id="img0",
            image=img,
            source="embedded",
            bbox=BoundingBox(x=50, y=50, width=200, height=100),
            width_px=200,
            height_px=100,
        )

        result = assemble_slide(page, layout, (), image_results=(extracted,))

        assert len(result.image_blocks) == 1
        ib = result.image_blocks[0]
        assert ib.region_id == "img0"
        assert ib.source == "embedded"
        assert ib.bbox.x == 50

    def test_full_page_image_from_page(self) -> None:
        """full_page_image comes from PageImage."""
        page = _make_page_image()
        layout = _make_layout_result()

        result = assemble_slide(page, layout, ())

        assert result.full_page_image is page.image

    def test_orphan_ocr_result_skipped(self) -> None:
        """OCR result without matching LayoutRegion is skipped."""
        page = _make_page_image()
        layout = _make_layout_result(regions=())
        orphan_ocr = _make_ocr_result(region_id="nonexistent")

        result = assemble_slide(page, layout, (orphan_ocr,))

        assert result.text_blocks == ()

    def test_formula_preserved(self) -> None:
        """Formula fields are preserved in TextBlock."""
        page = _make_page_image()
        region = _make_layout_region(region_id="r1")
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result(region_id="r1", is_formula=True, formula_latex="E=mc^2")

        result = assemble_slide(page, layout, (ocr,))

        tb = result.text_blocks[0]
        assert tb.is_formula is True
        assert tb.formula_latex == "E=mc^2"

    def test_uses_font_matches_when_provided(self) -> None:
        """Font matches override fallback when provided."""
        page = _make_page_image()
        region = _make_layout_region(region_id="r1", label=RegionLabel.TITLE)
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result(region_id="r1")

        font = FontMatch(
            region_id="r1",
            label=RegionLabel.TITLE,
            font_name="Google Sans",
            font_path=Path("/fonts/GoogleSans-Bold.ttf"),
            system_fallback="Arial",
            is_fallback=False,
        )

        result = assemble_slide(page, layout, (ocr,), font_matches=(font,))

        assert result.text_blocks[0].font_match.font_name == "Google Sans"
        assert result.text_blocks[0].font_match.is_fallback is False

    def test_falls_back_when_no_font_match(self) -> None:
        """Uses fallback when font_matches is empty."""
        page = _make_page_image()
        region = _make_layout_region(region_id="r1")
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result(region_id="r1")

        result = assemble_slide(page, layout, (ocr,), font_matches=())

        assert result.text_blocks[0].font_match.is_fallback is True
        assert result.text_blocks[0].font_match.font_name == "Arial"

    def test_uses_style_results_when_provided(self) -> None:
        """Style results are attached to TextBlocks."""
        page = _make_page_image()
        region = _make_layout_region(region_id="r1", label=RegionLabel.TITLE)
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result(region_id="r1")

        style = TextStyle(
            region_id="r1",
            font_size_pt=18,
            font_color_rgb=(128, 0, 0),
        )

        result = assemble_slide(page, layout, (ocr,), style_results=(style,))

        assert result.text_blocks[0].style is not None
        assert result.text_blocks[0].style.font_size_pt == 18
        assert result.text_blocks[0].style.font_color_rgb == (128, 0, 0)

    def test_style_none_when_not_provided(self) -> None:
        """TextBlock style is None when no style_results."""
        page = _make_page_image()
        region = _make_layout_region(region_id="r1")
        layout = _make_layout_result(regions=(region,))
        ocr = _make_ocr_result(region_id="r1")

        result = assemble_slide(page, layout, (ocr,))

        assert result.text_blocks[0].style is None


class TestAddTextBox:
    """Tests for _add_text_box()."""

    def _make_slide(self) -> tuple[Presentation, object]:
        """Create a minimal presentation with one blank slide."""
        prs = Presentation()
        prs.slide_width = Emu(int(13.333 * 914400))
        prs.slide_height = Emu(int(7.5 * 914400))
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        return prs, slide

    def test_adds_shape_to_slide(self) -> None:
        """Text box is added to the slide."""
        _, slide = self._make_slide()
        tb = _make_text_block()

        _add_text_box(slide, tb, dpi=300)

        assert len(slide.shapes) == 1

    def test_correct_position(self) -> None:
        """Text box position matches bbox at given DPI."""
        _, slide = self._make_slide()
        tb = _make_text_block(x=100, y=200, width=500, height=60)

        _add_text_box(slide, tb, dpi=300)

        shape = slide.shapes[0]
        emu_per_px = 914400 / 300
        assert shape.left == int(100 * emu_per_px)
        assert shape.top == int(200 * emu_per_px)
        assert shape.width == int(500 * emu_per_px)
        assert shape.height == int(60 * emu_per_px)

    def test_text_content(self) -> None:
        """Text box contains the correct text."""
        _, slide = self._make_slide()
        tb = _make_text_block(text="Hello World")

        _add_text_box(slide, tb, dpi=300)

        shape = slide.shapes[0]
        assert shape.text_frame.text == "Hello World"

    def test_font_size_matches_estimate(self) -> None:
        """Font size matches _estimate_font_size for bbox height."""
        _, slide = self._make_slide()
        tb = _make_text_block(height=60)

        _add_text_box(slide, tb, dpi=300)

        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        expected_size = Pt(_estimate_font_size(60, 300))
        assert run.font.size == expected_size

    def test_font_name_uses_font_match(self) -> None:
        """Font name comes from TextBlock's FontMatch (Arial for fallback)."""
        _, slide = self._make_slide()
        tb = _make_text_block()

        _add_text_box(slide, tb, dpi=300)

        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Arial"  # fallback font_match uses Arial

    def test_style_font_size_overrides_estimation(self) -> None:
        """When TextStyle is present, its font_size_pt is used."""
        _, slide = self._make_slide()
        tb = _make_text_block(height=60)
        # Attach a style with different font size
        styled = TextBlock(
            region_id=tb.region_id,
            bbox=tb.bbox,
            text=tb.text,
            font_match=tb.font_match,
            is_formula=tb.is_formula,
            style=TextStyle(
                region_id=tb.region_id,
                font_size_pt=24,
                font_color_rgb=(0, 0, 0),
            ),
        )

        _add_text_box(slide, styled, dpi=300)

        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(24)

    def test_style_font_color_applied(self) -> None:
        """When TextStyle is present, its font_color_rgb is applied."""
        _, slide = self._make_slide()
        tb = _make_text_block()
        styled = TextBlock(
            region_id=tb.region_id,
            bbox=tb.bbox,
            text=tb.text,
            font_match=tb.font_match,
            is_formula=tb.is_formula,
            style=TextStyle(
                region_id=tb.region_id,
                font_size_pt=11,
                font_color_rgb=(255, 0, 0),
            ),
        )

        _add_text_box(slide, styled, dpi=300)

        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(255, 0, 0)

    def test_title_alignment_is_center(self) -> None:
        """Title text blocks are center-aligned."""
        _, slide = self._make_slide()
        tb = _make_text_block(label=RegionLabel.TITLE)

        _add_text_box(slide, tb, dpi=300)

        from pptx.enum.text import PP_ALIGN

        shape = slide.shapes[0]
        assert shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

    def test_body_text_alignment_is_left(self) -> None:
        """Body text blocks are left-aligned."""
        _, slide = self._make_slide()
        tb = _make_text_block(label=RegionLabel.BODY_TEXT)

        _add_text_box(slide, tb, dpi=300)

        from pptx.enum.text import PP_ALIGN

        shape = slide.shapes[0]
        assert shape.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT


class TestBuildEditablePptx:
    """Tests for build_editable_pptx()."""

    def test_creates_valid_pptx(self, tmp_path: Path) -> None:
        """Output PPTX file exists and is non-empty."""
        slides = (_make_slide_content(),)
        output = tmp_path / "editable.pptx"

        result = build_editable_pptx(slides, output, dpi=300)

        assert result == output
        assert output.exists()
        assert output.stat().st_size > 0

    def test_correct_slide_count(self, tmp_path: Path) -> None:
        """PPTX has the correct number of slides."""
        slides = (
            _make_slide_content(page_number=0),
            _make_slide_content(page_number=1),
            _make_slide_content(page_number=2),
        )
        output = tmp_path / "multi.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        assert len(prs.slides) == 3

    def test_slides_contain_background_images(self, tmp_path: Path) -> None:
        """Each slide has a background image (full-slide picture)."""
        slides = (_make_slide_content(),)
        output = tmp_path / "bg.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        for slide in prs.slides:
            image_shapes = [
                s for s in slide.shapes if s.shape_type == 13  # PICTURE
            ]
            assert len(image_shapes) >= 1

    def test_text_boxes_present(self, tmp_path: Path) -> None:
        """Slides with text blocks have text box shapes."""
        tb = _make_text_block(text="Editable Text")
        slides = (_make_slide_content(text_blocks=(tb,)),)
        output = tmp_path / "text.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # Should have 1 background image + 1 text box = 2 shapes
        assert len(slide.shapes) == 2
        # At least one shape should be a text box (MSO_SHAPE_TYPE.AUTO_SHAPE = 1
        # or MSO_SHAPE_TYPE.PLACEHOLDER = 14, text boxes are typically shape_type=17)
        text_frames = [s for s in slide.shapes if s.has_text_frame]
        assert len(text_frames) >= 1

    def test_fallback_page_no_text_boxes(self, tmp_path: Path) -> None:
        """Fallback pages only have background image, no text boxes."""
        tb = _make_text_block(text="Should not appear")
        slides = (_make_slide_content(status="fallback", text_blocks=(tb,)),)
        output = tmp_path / "fallback.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # Only the background image, no text boxes
        assert len(slide.shapes) == 1

    def test_returns_output_path(self, tmp_path: Path) -> None:
        """Returns the output path for chaining."""
        output = tmp_path / "chain.pptx"

        result = build_editable_pptx((), output, dpi=300)

        assert result == output

    def test_empty_pages_creates_empty_pptx(self, tmp_path: Path) -> None:
        """Empty input creates a PPTX with zero slides."""
        output = tmp_path / "empty.pptx"

        build_editable_pptx((), output, dpi=300)

        prs = Presentation(str(output))
        assert len(prs.slides) == 0

    def test_output_parent_created(self, tmp_path: Path) -> None:
        """Parent directories are created if they don't exist."""
        output = tmp_path / "nested" / "dir" / "output.pptx"
        slides = (_make_slide_content(),)

        build_editable_pptx(slides, output, dpi=300)

        assert output.exists()

    def test_mixed_success_and_fallback(self, tmp_path: Path) -> None:
        """Mix of success and fallback pages works correctly."""
        tb = _make_text_block(text="Text")
        success = _make_slide_content(page_number=0, text_blocks=(tb,), status="success")
        fallback = _make_slide_content(page_number=1, status="fallback")
        output = tmp_path / "mixed.pptx"

        build_editable_pptx((success, fallback), output, dpi=300)

        prs = Presentation(str(output))
        assert len(prs.slides) == 2
        # First slide: background + text box
        assert len(prs.slides[0].shapes) == 2
        # Second slide: background only
        assert len(prs.slides[1].shapes) == 1

    def test_uses_background_image_when_available(self, tmp_path: Path) -> None:
        """When background_image is set, it is used as slide background."""
        bg = np.ones((2250, 4000, 3), dtype=np.uint8) * 200
        slide = _make_slide_content(background_image=bg)
        output = tmp_path / "with_bg.pptx"

        build_editable_pptx((slide,), output, dpi=300)

        assert output.exists()
        prs = Presentation(str(output))
        assert len(prs.slides) == 1
        # Should have a picture shape (background)
        assert len(prs.slides[0].shapes) >= 1

    def test_falls_back_to_full_page_image(self, tmp_path: Path) -> None:
        """When background_image is None, full_page_image is used."""
        slide = _make_slide_content()
        assert slide.background_image is None
        output = tmp_path / "fallback_bg.pptx"

        build_editable_pptx((slide,), output, dpi=300)

        assert output.exists()
        prs = Presentation(str(output))
        assert len(prs.slides) == 1

    def test_image_blocks_placed_as_pictures(self, tmp_path: Path) -> None:
        """Image blocks are placed as picture shapes in PPTX."""
        img = np.zeros((100, 200, 3), dtype=np.uint8)
        ib = ImageBlock(
            region_id="img0",
            bbox=BoundingBox(x=100, y=200, width=500, height=300),
            image=img,
            source="cropped",
        )
        slides = (_make_slide_content(image_blocks=(ib,)),)
        output = tmp_path / "images.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # background + 1 image block = 2 shapes
        assert len(slide.shapes) == 2

    def test_image_block_position(self, tmp_path: Path) -> None:
        """Image block position matches bbox at given DPI."""
        img = np.zeros((100, 200, 3), dtype=np.uint8)
        ib = ImageBlock(
            region_id="img0",
            bbox=BoundingBox(x=100, y=200, width=500, height=300),
            image=img,
            source="cropped",
        )
        slides = (_make_slide_content(image_blocks=(ib,)),)
        output = tmp_path / "img_pos.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # Second shape is the image block (first is background)
        img_shape = [s for s in slide.shapes if s.shape_type == 13][1]
        emu_per_px = 914400 / 300
        assert img_shape.left == int(100 * emu_per_px)
        assert img_shape.top == int(200 * emu_per_px)
        assert img_shape.width == int(500 * emu_per_px)
        assert img_shape.height == int(300 * emu_per_px)

    def test_fallback_page_no_image_blocks(self, tmp_path: Path) -> None:
        """Fallback pages don't place image blocks."""
        img = np.zeros((100, 200, 3), dtype=np.uint8)
        ib = ImageBlock(
            region_id="img0",
            bbox=BoundingBox(x=100, y=200, width=500, height=300),
            image=img,
            source="cropped",
        )
        slides = (_make_slide_content(status="fallback", image_blocks=(ib,)),)
        output = tmp_path / "fallback_img.pptx"

        build_editable_pptx(slides, output, dpi=300)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # Only background, no image block
        assert len(slide.shapes) == 1
